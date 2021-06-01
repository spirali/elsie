from elsie import SlideDeck


def export_pptx(slide_deck: SlideDeck, filename: str, **render_args):
    import pptx

    pr = pptx.Presentation()

    for layout in pr.slide_layouts:
        if layout.name == "Blank":
            break
    else:
        raise Exception("Blank layout not found")

    images = slide_deck.render(None, export_type="png", **render_args)
    for image in images:
        slide = pr.slides.add_slide(layout)
        slide.shapes.add_picture(image, 0, 0, pr.slide_width, pr.slide_height)
    pr.save(filename)
